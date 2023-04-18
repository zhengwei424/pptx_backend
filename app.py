import os, sys
import json
import pptx
from flask import Flask, request, jsonify, render_template, send_from_directory
from flask_cors import CORS
from pptx import Presentation
from week_func import SetWeaklyReport
from tool_func import generate_table_data
from month_func import GetMonthlyReportsData, SetMonthlyReport
from customMonthlySummary_func import SetMonthlySummaryReport

ALLOWED_EXTENSIONS = {'pptx'}

app = Flask(__name__, static_folder="static", template_folder="template", static_url_path="")
app.config['UPLOAD_FOLDER'] = "data"
CORS(app, supports_credentials=True)


def allowed_file(filename):
    '''
    判断文件后缀是否满足上传需求
    :param filename: 文件名
    :return: True or False
    '''
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


# 获取周报文件列表
@app.route("/weeklyReports", methods=["GET"])
def weekly_reports():
    report_files = os.listdir(app.config["UPLOAD_FOLDER"] + "/weeklyReports")
    report_files.sort()
    result = []
    for file in report_files:
        item = dict({
            "name": str(file)
        })
        result.append(item)
    return jsonify(result)


# 获取周报json文件列表 / 更新周报JSON文件
@app.route("/weeklyReportsJson", methods=["GET", "POST"])
def weekly_reports_json():
    if request.method == "GET":
        report_files = os.listdir(app.config["UPLOAD_FOLDER"] + "/historyWeeklyData")
        report_files.sort()
        result = []
        for file in report_files:
            content = ''
            with open(app.config['UPLOAD_FOLDER'] + '/historyWeeklyData/' + file, mode='r', encoding='UTF-8') as f:
                while True:
                    line = f.readline()
                    if not line:
                        break
                    content += line
            item = dict({
                "fileName": str(file),
                "fileContent": json.dumps(json.loads(content), indent=2, ensure_ascii=False)
            })
            result.append(item)
        return jsonify(result)
    if request.method == "POST":
        data = request.json
        try:
            json.loads(data["fileContent"])
        except Exception as e:
            resp_data = {
                "code": 1,
                "msg": str(e)
            }
            return jsonify(resp_data)

        with open(app.config['UPLOAD_FOLDER'] + '/historyWeeklyData/' + data["fileName"], mode='w') as f:
            f.write(json.dumps(json.loads(data["fileContent"])))
        resp_data = {
            "code": 0,
            "msg": "JSON更新成功"
        }
        return jsonify(resp_data)


# 获取月报文件列表
@app.route("/monthlyReports", methods=["GET"])
def monthly_reports():
    report_files = os.listdir(app.config["UPLOAD_FOLDER"] + "/monthlyReports")
    report_files.sort()
    result = []
    for file in report_files:
        item = dict({
            "name": str(file)
        })
        result.append(item)
    return jsonify(result)


# 获取月报json文件列表 / 更新周报JSON文件
@app.route("/monthlyReportsJson", methods=["GET"])
def monthly_reports_json():
    report_files = os.listdir(app.config["UPLOAD_FOLDER"] + "/historyMonthlyData")
    report_files.sort()
    result = []
    for file in report_files:
        content = ''
        with open(app.config['UPLOAD_FOLDER'] + '/historyMonthlyData/' + file, mode='r', encoding='UTF-8') as f:
            while True:
                line = f.readline()
                if not line:
                    break
                content += line
        item = dict({
            "fileName": str(file)
        })
        result.append(item)
    return jsonify(result)


# 获取月报汇总文件列表
@app.route("/monthlySummaryReports", methods=["GET"])
def monthly_summary_reports():
    report_files = os.listdir(app.config["UPLOAD_FOLDER"] + "/monthlySummaryReports")
    report_files.sort()
    result = []
    for file in report_files:
        item = dict({
            "name": str(file)
        })
        result.append(item)
    return jsonify(result)


# 上传周报
@app.route('/weeklyReports/upload', methods=['GET', 'POST'])
def weekly_reports_upload():
    file = request.files.get("file")
    # extra_args = dict(request.form)
    resp_data = {}
    if allowed_file(file.filename):
        if os.path.exists(app.config["UPLOAD_FOLDER"] + "/weeklyReports/" + file.filename):
            resp_data = {
                "msg": file.filename + "文件已存在!"
            }
        else:
            file.save(app.config["UPLOAD_FOLDER"] + "/weeklyReports/" + file.filename)
            resp_data = {
                "msg": file.filename + "文件上传成功!"
            }
    else:
        resp_data = {
            "msg": "只能上传后缀为.pptx的文件"
        }

    return jsonify(resp_data)


# 上传月报
@app.route('/monthlyReports/upload', methods=['GET', 'POST'])
def monthly_reports_upload():
    file = request.files.get("file")
    # extra_args = dict(request.form)
    resp_data = {}
    if allowed_file(file.filename):
        if os.path.exists(app.config["UPLOAD_FOLDER"] + "/monthlyReports/" + file.filename):
            resp_data = {
                "msg": file.filename + "文件已存在!"
            }
        else:
            file.save(app.config["UPLOAD_FOLDER"] + "/monthlyReports/" + file.filename)
            resp_data = {
                "msg": file.filename + "文件上传成功!"
            }
    else:
        resp_data = {
            "msg": "只能上传后缀为.pptx的文件"
        }

    return jsonify(resp_data)


# 下载周报
@app.route('/weeklyReports/download/<filename>', methods=["GET"])
def weekly_reports_download(filename):
    # send_static_file会在static目录下寻找文件
    # return app.send_static_file("weeklyReports/" + filename)
    return send_from_directory(directory=app.config['UPLOAD_FOLDER'] + "/weeklyReports/", path=filename)


# 下载月报
@app.route('/monthlyReports/download/<filename>', methods=["GET"])
def monthly_reports_download(filename):
    # send_static_file会在static目录下寻找文件
    # return app.send_static_file("monthlyReports/" + filename)
    return send_from_directory(directory=app.config['UPLOAD_FOLDER'] + "/monthlyReports/", path=filename)


# 下载月报汇总
@app.route('/monthlySummaryReports/download/<filename>', methods=["GET"])
def monthly_summary_reports_download(filename):
    # send_static_file会在static目录下寻找文件
    # return app.send_static_file("monthlyReports/" + filename)
    return send_from_directory(directory=app.config['UPLOAD_FOLDER'] + "/monthlySummaryReports/", path=filename)


# 生成周报
@app.route("/weeklyReportsData", methods=["POST"])
def generate_weekly_report():
    # request.json是一个字典，接收post请求的data数据
    # {"weeklyData": weeklyData, "formdata": this.formData}
    post_data = request.json
    year = post_data.get("formdata").get("year")
    month = post_data.get("formdata").get("month")
    week = post_data.get("formdata").get("week")
    # 通过周报的json数据重新生成周报
    weekly_report_json_data = post_data.get("formdata").get("weekly_report_json_data")

    # 响应消息
    msg = ""
    history_weekly_data_file_name = year + month + week + ".json"
    if post_data.get("status") == 0:
        if os.path.exists(app.config["UPLOAD_FOLDER"] + "/historyWeeklyData/" + history_weekly_data_file_name):
            msg = "历史周报JSON文件已存在"
            return jsonify({
                "code": 1,
                "msg": msg
            })

        with open(app.config["UPLOAD_FOLDER"] + "/historyWeeklyData/" + history_weekly_data_file_name, mode="x",
                  encoding="utf-8") as f:
            f.write(json.dumps(post_data))

    try:
        # inspect_data
        inspect_data = []
        for i in post_data.get("weeklyData").get("inspect")[0].items():
            if i[0] != "index":
                inspect_data.append(i[1])

        # change_data
        change_data = []
        for i in post_data.get("weeklyData").get("change"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
            change_data.append(tmp)

        # release_data
        release_data = []
        for i in post_data.get("weeklyData").get("release"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
            release_data.append(tmp)

        # permissionManagement_data(有序号)
        permission_management_data = []
        for i in post_data.get("weeklyData").get("permissionManagement"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
                    else:
                        tmp.insert(0, str(int(j[1]) + 1))
            permission_management_data.append(tmp)

        # cooperation_data(有序号)
        cooperation_data = []
        for i in post_data.get("weeklyData").get("cooperation"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
                    else:
                        tmp.insert(0, str(int(j[1]) + 1))
            cooperation_data.append(tmp)

        # problem_data
        problem_data = []
        for i in post_data.get("weeklyData").get("problem"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
            problem_data.append(tmp)

        # workingPlan_data(有序号)
        working_plan_data = []
        for i in post_data.get("weeklyData").get("workingPlan"):
            tmp = []
            if i.items():
                for j in i.items():
                    if j[0] != "index":
                        tmp.append(j[1])
                    else:
                        tmp.insert(0, str(int(j[1]) + 1))
            working_plan_data.append(tmp)

        # events_count
        events_count = [
            len(change_data),  # 变更
            len(permission_management_data),  # 资源权限管理
            len(cooperation_data),  # 配合操作
            len(release_data),  # 支撑发版
            len(problem_data),  # 问题和告警处理
            0  # 故障处理
        ]

        # 获取运行情况分析ppt中所需所需数据
        cluster_pie_data, cluster_table_data = generate_table_data()

        # 编写ppt
        # 这种打开方式适合ppt2007及最新，不适合ppt2003及以前。支持stringio/bytesio stream
        prs = Presentation(
            app.config["UPLOAD_FOLDER"] + "/template/week.pptx")  # type: pptx.presentation.Presentation # 设置type，会有代码提示
        wr = SetWeaklyReport(prs)
        # 运维工作统计
        wr.slide_1(events_count)

        # 巡检
        wr.slide_2(inspect_data)

        # 变更
        wr.slide_3(change_data)

        # 支撑发版
        wr.slide_4(release_data)

        # 资源权限管理
        wr.slide_5(permission_management_data)

        # 配合操作
        wr.slide_6(cooperation_data)

        # 问题及告警
        wr.slide_7(problem_data)

        # 运行情况分析
        wr.slide_8(cluster_pie_data, cluster_table_data)

        # 下周工作计划
        wr.slide_9(working_plan_data)

        # 保存pptx
        weeklyReportFileName = year + month + week + ".pptx"

        # if os.path.exists(app.config["UPLOAD_FOLDER"] + "/weeklyReports/" + weeklyReportFileName):
        #     msg = "周报已存在"
        #     return jsonify({
        #         "code": 1,
        #         "msg": msg
        #     })
        # else:
        prs.save(app.config["UPLOAD_FOLDER"] + "/weeklyReports/" + weeklyReportFileName)
        return jsonify({
            "code": 0,
            "msg": "周报生成/更新成功"
        })
    except Exception as e:
        msg = "周报生成/更新失败" + str(e)
        return jsonify({
            "code": 1,
            "msg": msg
        })


# 生成月报
@app.route("/monthlyReportsData", methods=["POST"])
def generate_monthly_report():
    # 获取周报文件
    files = request.json

    # 月报数据初始化
    # 巡检[巡检次数、提交报告次数、正常次数、异常次数]
    month_inspect_data = [0, 0, 0, 0]
    # 运维工作统计[变更、资源权限管理、配合操作、支撑发版、故障及问题处理]
    month_event_count = [0, 0, 0, 0, 0]
    work_summary = ["本月主要工作："]
    # 变更内容
    month_change_data = []
    # 资源权限管理
    month_permission_management_data = []
    # 配合操作
    month_cooperation_data = []
    # 支撑发版
    month_release_data = []
    # 问题处理
    month_problem_data = []
    # 运行情况分析
    month_analyse_data = []

    try:
        for file in files:
            prs = pptx.Presentation(app.config["UPLOAD_FOLDER"] + "/weeklyReports/" + file)
            mr = GetMonthlyReportsData(prs)
            event_count = mr.get_event_count()
            inspect_data = mr.get_inspect_data()
            change_data = mr.get_change_data()
            release_data = mr.get_release_data()
            permission_management_data = mr.get_permission_management_data()
            cooperation_data = mr.get_cooperation_data()
            problem_data = mr.get_problem_data()
            analyse_data = mr.get_analyse()

            month_inspect_data[0] += inspect_data[0]  # 巡检次数
            month_inspect_data[1] += inspect_data[2]  # 提交报告次数
            month_inspect_data[2] += inspect_data[0] - inspect_data[1]  # 正常次数
            month_inspect_data[3] += inspect_data[1]  # 异常次数

            month_event_count[0] += event_count[0]  # 变更
            month_event_count[1] += event_count[1]  # 资源权限管理
            month_event_count[2] += event_count[2]  # 配合操作
            month_event_count[3] += event_count[3]  # 支撑发版
            month_event_count[4] += event_count[4] + event_count[5]  # 故障及问题

            if change_data:
                for item in change_data:
                    month_change_data.append(item[2])

            for item in permission_management_data:
                month_permission_management_data.append(item)

            for item in cooperation_data:
                month_cooperation_data.append(item)

            for item in release_data:
                month_release_data.append(item)

            for item in problem_data:
                month_problem_data.append(item)

            month_analyse_data.append(analyse_data)

        if month_change_data:
            work_summary += month_change_data
        if month_cooperation_data:
            work_summary.append("日常配合应用运维或研发人员做问题排查" + str(len(month_cooperation_data)) + "次")
        if month_release_data:
            work_summary.append("支撑发版" + str(len(month_release_data)) + "次")
        if month_problem_data:
            work_summary.append("联通云平台问题处理" + str(len(month_problem_data)) + "次")

        # 月报JSON数据
        monthly_report_json_data = {
            # 巡检[巡检次数、提交报告次数、正常次数、异常次数]
            "month_inspect_data": month_inspect_data,
            # 运维工作统计[变更、资源权限管理、配合操作、支撑发版、故障及问题处理]
            "month_event_count": month_event_count,
            # 变更内容
            "month_change_data": month_change_data,
            # 资源权限管理
            "month_permission_management_data": month_permission_management_data,
            # 配合操作
            "month_cooperation_data": month_cooperation_data,
            # 支撑发版
            "month_release_data": month_release_data,
            # 问题处理
            "month_problem_data": month_problem_data,
            # 运行情况分析
            "month_analyse_data": month_analyse_data
        }

        history_monthly_data_file_name = app.config["UPLOAD_FOLDER"] + "/historyMonthlyData/" + files[-1].split(".")[0][
                                                                                                :-2] + ".json"

        if os.path.exists(history_monthly_data_file_name):
            return jsonify({
                "code": 1,
                "msg": "月报JSON文件已存在"
            })
        with open(history_monthly_data_file_name, mode="x",
                  encoding="utf-8") as f:
            f.write(json.dumps(monthly_report_json_data))

        # 资源权限管理添加序号
        i = 1
        for item in month_permission_management_data:
            item.insert(0, str(i))
            i += 1

        # 配合操作添加序号
        i = 1
        for item in month_cooperation_data:
            item.insert(0, str(i))
            i += 1

    except Exception as e:
        return jsonify({
            "code": 1,
            "msg": "获取周报数据失败" + str(e)
        })

    try:
        # 合成月报
        prs = pptx.Presentation(app.config["UPLOAD_FOLDER"] + "/template/month.pptx")
        mr = SetMonthlyReport(prs)

        # 巡检
        mr.slide_1(month_inspect_data)

        # 运维工作统计
        mr.slide_2(month_event_count, work_summary)

        # 变更
        mr.slide_3(month_change_data)

        # 资源权限管理(有序号)
        mr.slide_4(month_permission_management_data)

        # 配合操作(有序号)
        mr.slide_5(month_cooperation_data)

        # 支撑发版
        mr.slide_6(month_release_data)

        # 问题处理
        mr.slide_7(month_problem_data)

        # 运行情况分析
        mr.slide_8(month_analyse_data)

        # 保存月报
        # if os.path.exists(app.config["UPLOAD_FOLDER"] + "/monthlyReports/" + files[-1].split(".")[0][:-2] + ".pptx"):
        #     msg = "月报已存在"
        #     return jsonify({
        #         "code": 1,
        #         "msg": msg
        #     })
        # else:
        prs.save(app.config["UPLOAD_FOLDER"] + "/monthlyReports/" + files[-1].split(".")[0][:-2] + ".pptx")

        return jsonify({
            "code": 0,
            "msg": "月报生成成功"
        })
    except Exception as e:
        return jsonify({
            "code": 1,
            "msg": "月报生成失败" + str(e)
        })


# 月报汇总
@app.route("/monthlySummaryData", methods=["POST"])
def generate_monthly_summary_report():
    # 获取周报文件
    files = request.json

    # 月报汇总范围
    months = []
    # 月报数据初始化
    # 巡检[巡检次数、提交报告次数、正常次数、异常次数]
    monthly_summary_inspect_data = [0, 0, 0, 0]
    # 运维工作统计[变更、资源权限管理、配合操作、支撑发版、故障及问题处理]
    monthly_summary_event_count = [0, 0, 0, 0, 0]
    work_summary = [files[0].split(".")[0][4:] + "-" + files[-1].split(".")[0][4:] + "月主要工作："]
    # 变更内容
    monthly_summary_change_data = []
    # 资源权限管理
    monthly_summary_permission_management_data = []
    # 配合操作
    monthly_summary_cooperation_data = []
    # 支撑发版
    monthly_summary_release_data = []
    # 问题处理
    monthly_summary_problem_data = []
    # 运行情况分析
    monthly_summary_analyse_data = []
    try:
        for file in files:
            file_json_data = {}
            with open(app.config["UPLOAD_FOLDER"] + "/historyMonthlyData/" + file, mode='r', encoding="UTF-8") as f:
                line = f.readline()
                if line:
                    file_json_data = json.loads(line)

            monthly_summary_inspect_data[0] += file_json_data.get("month_inspect_data")[0]  # 巡检次数
            monthly_summary_inspect_data[1] += file_json_data.get("month_inspect_data")[1]  # 提交报告次数
            monthly_summary_inspect_data[2] += file_json_data.get("month_inspect_data")[2]  # 正常次数
            monthly_summary_inspect_data[3] += file_json_data.get("month_inspect_data")[3]  # 异常次数

            monthly_summary_event_count[0] += file_json_data.get("month_event_count")[0]  # 变更
            monthly_summary_event_count[1] += file_json_data.get("month_event_count")[1]  # 资源权限管理
            monthly_summary_event_count[2] += file_json_data.get("month_event_count")[2]  # 配合操作
            monthly_summary_event_count[3] += file_json_data.get("month_event_count")[3]  # 支撑发版
            monthly_summary_event_count[4] += file_json_data.get("month_event_count")[4]  # 故障及问题ch处理

            monthly_summary_change_data += file_json_data.get("month_change_data")

            monthly_summary_permission_management_data += file_json_data.get("month_permission_management_data")

            monthly_summary_cooperation_data += file_json_data.get("month_cooperation_data")

            monthly_summary_release_data += file_json_data.get("month_release_data")

            monthly_summary_problem_data += file_json_data.get("month_problem_data")

            monthly_summary_analyse_data.append(file_json_data.get("month_analyse_data")[-1])

            months.append(files[-1].split(".")[0][4:] + "月")

        if monthly_summary_change_data:
            work_summary += monthly_summary_change_data
        if monthly_summary_cooperation_data:
            work_summary.append(
                "日常配合应用运维或研发人员做问题排查" + str(len(monthly_summary_cooperation_data)) + "次")
        if monthly_summary_release_data:
            work_summary.append("支撑发版" + str(len(monthly_summary_release_data)) + "次")
        if monthly_summary_problem_data:
            work_summary.append("联通云平台问题处理" + str(len(monthly_summary_problem_data)) + "次")

        # 资源权限管理添加序号
        i = 1
        for item in monthly_summary_permission_management_data:
            item.insert(0, str(i))
            i += 1

        # 配合操作添加序号
        i = 1
        for item in monthly_summary_cooperation_data:
            item.insert(0, str(i))
            i += 1

    except Exception as e:
        return jsonify({
            "code": 1,
            "msg": "获取月报JSON报数据失败" + str(e)
        })

    try:
        # 汇总月报
        prs = pptx.Presentation(app.config["UPLOAD_FOLDER"] + "/template/monthsSummary.pptx")
        mr = SetMonthlySummaryReport(prs)

        # 巡检
        mr.slide_1(monthly_summary_inspect_data)

        # 运维工作统计
        mr.slide_2(monthly_summary_event_count, work_summary)

        # 变更
        mr.slide_3(monthly_summary_change_data)

        # 资源权限管理(有序号)
        mr.slide_4(monthly_summary_permission_management_data)

        # 配合操作(有序号)
        mr.slide_5(monthly_summary_cooperation_data)

        # 支撑发版
        mr.slide_6(monthly_summary_release_data)

        # 问题处理
        mr.slide_7(monthly_summary_problem_data)

        # 运行情况分析
        mr.slide_8(monthly_summary_analyse_data, months)

        # 保存月报
        # if os.path.exists(app.config["UPLOAD_FOLDER"] + "/monthlyReports/" + files[-1].split(".")[0][:-2] + ".pptx"):
        #     msg = "月报已存在"
        #     return jsonify({
        #         "code": 1,
        #         "msg": msg
        #     })
        # else:
        monthly_summary_report_file_name = app.config["UPLOAD_FOLDER"] + "/monthlySummaryReports/" + \
                                           files[0].split(".")[0][:4] + "年" + \
                                           files[0].split(".")[0][4:] + "-" + files[-1].split(".")[0][4:] + \
                                           "月月报汇总.pptx"
        prs.save(monthly_summary_report_file_name)

        return jsonify({
            "code": 0,
            "msg": "月报汇总成功"
        })
    except Exception as e:
        return jsonify({
            "code": 1,
            "msg": "月报汇总失败" + str(e)
        })


@app.route("/")
def index():
    return render_template("index.html")


if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000, debug=False)
