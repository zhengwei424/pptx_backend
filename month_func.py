import pptx.presentation
import pptx.slide
import pptx.shapes.base
import pptx.table
import pptx.chart.chart
import pptx.chart.plot
import pptx.chart.series
import pptx.shapes.graphfrm
import pptx.dml.chtfmt
import pptx.dml.fill
import pptx.dml.line
import pptx.chart.axis
import pptx.shapes.placeholder
import pptx.shapes.autoshape
from tool_func import set_text_frame_paragraph_format, set_cell_format, create_line_chart, create_column_chart
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm


# presentation -> slide -> shapes -> placeholder,graphfrm -> chart(table -> cell) -> text_frame -> paragraphs -> font
#            |      +-> placeholder                            | - > plot(感觉是扇形的扇子，柱形的柱子) -> categories
#            |                                                          | -> series
#            |->slide_master -> slide_layout
# 运维工作统计
class GetMonthlyReportsData(object):
    def __init__(self, prs):
        self._prs = prs  # type: pptx.presentation.Presentation

    # 运维工作统计
    def get_event_count(self):
        # 返回运维工作统计结果
        events_count = []
        # 获取slide
        slide = self._prs.slides[8]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                events_count = [
                    int(tb.cell(3, 1).text),  # 变更
                    int(tb.cell(3, 2).text),  # 资源权限管理
                    int(tb.cell(3, 3).text),  # 配合操作
                    int(tb.cell(3, 4).text),  # 支撑发版
                    int(tb.cell(3, 5).text),  # 问题和告警
                    int(tb.cell(3, 6).text),  # 故障
                ]
        return events_count

    # 巡检
    def get_inspect_data(self):
        # 返回巡检数据
        inspect_data = []
        # 获取slide
        slide = self._prs.slides[9]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                inspect_data = [
                    int(tb.cell(3, 1).text),  # 巡检次数
                    int(tb.cell(3, 2).text),  # 异常次数
                    int(tb.cell(3, 3).text)  # 提交报告次数
                ]
        return inspect_data

    # 变更
    def get_change_data(self):
        # 返回变更数据
        change_data = []
        # 获取slide
        slide = self._prs.slides[10]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(1, len(tb.rows)):
                    item = []
                    if not tb.cell(row_idx, 1).text:
                        break
                    for col_idx in range(len(tb.columns)):
                        item.append(tb.cell(row_idx, col_idx).text)
                    change_data.append(item)
        return change_data

    # 支撑发版
    def get_release_data(self):
        # 返回支撑发版数据
        release_data = []
        # 获取slide
        slide = self._prs.slides[11]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(1, len(tb.rows)):
                    item = []
                    if not tb.cell(row_idx, 1).text:
                        break
                    for col_idx in range(len(tb.columns)):
                        item.append(tb.cell(row_idx, col_idx).text)
                    release_data.append(item)
        return release_data

    # 资源权限管理(有序号)
    def get_permission_management_data(self):
        # 返回资源权限管理数据
        permission_management_data = []
        # 获取slide
        slide = self._prs.slides[12]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(1, len(tb.rows)):
                    item = []
                    if not tb.cell(row_idx, 1).text:
                        break
                    for col_idx in range(1, len(tb.columns)):
                        item.append(tb.cell(row_idx, col_idx).text)
                    permission_management_data.append(item)
        return permission_management_data

    # 配合操作(有序号)
    def get_cooperation_data(self):
        # 返回配合操作数据
        cooperation_data = []
        # 获取slide
        slide = self._prs.slides[13]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(1, len(tb.rows)):
                    item = []
                    if not tb.cell(row_idx, 1).text:
                        break
                    for col_idx in range(1, len(tb.columns)):
                        item.append(tb.cell(row_idx, col_idx).text)
                    cooperation_data.append(item)
        return cooperation_data

    # 问题及告警
    def get_problem_data(self):
        # 返回问题与告警数据
        problem_data = []
        # 获取slide
        slide = self._prs.slides[14]  # type:pptx.slide.Slide
        # 获取shape
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(1, len(tb.rows)):
                    item = []
                    if not tb.cell(row_idx, 1).text:
                        break
                    for col_idx in range(len(tb.columns)):
                        item.append(tb.cell(row_idx, col_idx).text)
                    problem_data.append(item)
        return problem_data

    # 运行情况分析
    def get_analyse(self):
        # 运行分析数据，数据结构是{"<cluster_name>": {"cpu": [], "memory": [], "table_data": []}}
        analyse_data = {}
        cluster_name = ""
        # 第一张运行情况分析slide的索引值
        slide_idx = 18
        # cluster names
        cluster_names = [
            "FCP业务集群",
            # "内网微服务区",
            # "内网后台区",
            # "外网微服务区",
            # "外网后台区"
        ]
        for item in cluster_names:
            analyse_data[item] = {"cpu": [], "memory": [], "table_data": [], "summary": ""}
        # cluster name
        cluster_name = ""
        while True:
            # 获取slide
            slide = self._prs.slides[slide_idx]  # type:pptx.slide.Slide
            # 获取shape
            shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
            if "运行情况分析" not in shape.text:
                break
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart  # type: pptx.chart.chart.Chart
                    plot = chart.plots[0]  # type: pptx.chart.plot._BasePlot
                    # categories = plot.categories
                    # categories_labels = [c.label for c in categories]
                    # print(categories_labels)
                    cluster_name = chart.chart_title.text_frame.paragraphs[0].text.split(" ")[0]

                    series = plot.series
                    for item in series:
                        if item.name == "cpu":
                            analyse_data[cluster_name]["cpu"] = list(item.values)
                        if item.name == "memory":
                            analyse_data[cluster_name]["memory"] = list(item.values)
                if shape.has_table:
                    tb = shape.table  # type: pptx.table.Table
                    # start_idx = 0
                    # if tb.cell(0, 0).text == "项目名称":
                    #     start_idx = 1
                    for row_idx in range(len(tb.rows)):
                        tmp = []
                        for col_idx in range(len(tb.columns)):
                            tmp.append(tb.cell(row_idx, col_idx).text)
                        analyse_data[cluster_name]["table_data"].append(tmp)
                if isinstance(shape, pptx.shapes.autoshape.Shape) and cluster_name and cluster_name in shape.text:
                    analyse_data[cluster_name]["summary"] = shape.text
            slide_idx += 1
        for key in analyse_data.keys():
            tmp_text = ""
            for item in analyse_data[key]["table_data"]:
                if item[0] != tmp_text:
                    tmp_text = item[0]
                else:
                    item[0] = ""
        return analyse_data


class SetMonthlyReport(object):
    def __init__(self, prs):
        self._prs = prs  # type: pptx.presentation.Presentation
        self._current_slide_idx = 0  # 当前操作的幻灯片索引

    def insert_slide(self, insert_index: int):
        xml_slides = self._prs.slides._sldIdLst
        slideIdList = list(xml_slides)
        item = slideIdList.pop()
        xml_slides.insert(insert_index, item)

    # 巡检
    def slide_1(self, month_inspect_data: list):
        """巡检"""
        self._current_slide_idx = 5
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(1, 1),  # 巡检次数
                    tb.cell(1, 2),  # 提交报告次数
                    tb.cell(1, 3),  # 正常次数
                    tb.cell(1, 4),  # 异常次数
                ]
                for index in range(len(cells)):
                    cells[index].text = str(month_inspect_data[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=12,
                                                    align=PP_ALIGN.CENTER
                                                    )

    # 运维工作统计
    def slide_2(self, month_event_count: list, work_summary: list):
        """运维工作统计"""
        self._current_slide_idx = 6
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        # 填充表格数据
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(1, 1),  # 变更
                    tb.cell(1, 2),  # 资源及权限管理
                    tb.cell(1, 3),  # 配合操作
                    tb.cell(1, 4),  # 支撑发版
                    tb.cell(1, 5),  # 故障及问题处理
                ]
                for index in range(len(cells)):
                    cells[index].text = str(month_event_count[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=10,
                                                    align=PP_ALIGN.CENTER
                                                    )

        # 初始化柱状图数据
        create_column_chart(slide=slide,
                            categories=['变更', '资源权限管理', '配合操作', '支撑发版', '故障和问题处理'],
                            series={"运维工作统计": month_event_count},
                            position_and_size=[Cm(1.1), Cm(4.5), Cm(10), Cm(8.5)],
                            chart_title="本月运维工作统计")

        # 总结本月主要工作
        text_list = work_summary
        txbox = slide.shapes.add_textbox(Cm(13.5), Cm(5.5), Cm(11), Cm(5))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # 自动换行
        tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "\n●".join(text_list)
        set_text_frame_paragraph_format(tfp, font_size=10, font_bold=True, font_color="FF0000")

    # 变更
    def slide_3(self, month_change_data):
        """变更"""
        self._current_slide_idx = 7
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                for row_idx in range(len(tb.rows)):
                    if row_idx == len(month_change_data):
                        break
                    cell = tb.cell(row_idx, 1)
                    cell.text = month_change_data[row_idx]
                    set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                    font_size=10,
                                                    align=PP_ALIGN.LEFT)

    # 资源权限管理(有序号)
    def slide_4(self, month_permission_management_data: list):
        """资源权限管理"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(month_permission_management_data) == 0:
            pages = 1
        elif len(month_permission_management_data) % (row - 1):
            pages = len(month_permission_management_data) // (row - 1) + 1
        else:
            pages = len(month_permission_management_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-资源权限管理"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(3)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not month_permission_management_data:
                    break
                data = month_permission_management_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 配合操作(有序号)
    def slide_5(self, month_cooperation_data: list):
        """配合操作"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(month_cooperation_data) == 0:
            pages = 1
        elif len(month_cooperation_data) % (row - 1):
            pages = len(month_cooperation_data) // (row - 1) + 1
        else:
            pages = len(month_cooperation_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-配合操作及排障"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)
            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(3)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not month_cooperation_data:
                    break
                data = month_cooperation_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 支撑发版
    def slide_6(self, month_release_data: list):
        """支撑发版"""
        # 表格行数
        row = 15
        # 表格页数初始化
        pages = 0
        if len(month_release_data) == 0:
            pages = 1
        elif len(month_release_data) % (row - 1):
            pages = len(month_release_data) // (row - 1) + 1
        else:
            pages = len(month_release_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-支撑发版"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
            table.columns[0].width = Cm(3)
            table.columns[1].width = Cm(4)
            table.columns[2].width = Cm(3)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(5)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
            ]
            table_header = [
                "专业",
                "发版时间",
                "次数",
                "工作内容",
                "异常情况处理",
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not month_release_data:
                    break
                data = month_release_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 问题处理
    def slide_7(self, month_problem_data: list):
        """问题处理"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(month_problem_data) == 0:
            pages = 1
        elif len(month_problem_data) % (row - 1):
            pages = len(month_problem_data) // (row - 1) + 1
        else:
            pages = len(month_problem_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-问题处理"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(5)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(5)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
            ]
            table_header = [
                "专业",
                "问题详细描述",
                "处理结果",
                "原因分析",
                "后续建议",
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not month_problem_data:
                    break
                data = month_problem_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 10. 运行情况分析
    def slide_8(self, month_analyse_data: list):
        # 解析数据
        data_length = len(month_analyse_data)
        #cluster_names = ["FCP业务集群", "内网微服务区", "内网后台区", "外网微服务区", "外网后台区"]
        cluster_names = ["FCP业务集群"]
        weeks = ["第1周", "第2周", "第3周", "第4周", "第5周"]
        cpu_allocated = {}
        memory_allocated = {}
        cluster_table_data = {}
        cluster_summary = {}
        for item in cluster_names:
            cpu_allocated[item] = []
            memory_allocated[item] = []
            cluster_table_data[item] = []

        for index in range(data_length):
            for key in month_analyse_data[index].keys():
                cpu_allocated[key].append(month_analyse_data[index][key]["cpu"][0])
                memory_allocated[key].append(month_analyse_data[index][key]["memory"][0])
                if index == (data_length - 1):
                    cluster_table_data[key] = month_analyse_data[index][key]["table_data"]
                    cluster_summary[key] = month_analyse_data[index][key]["summary"]

        # 添加一张幻灯片
        self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

        # 当前slide
        slide = self._prs.slides[self._current_slide_idx + 6]  # type:pptx.slide.Slide

        # 幻灯片标题
        slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
        p = slidePlaceholder.text_frame.paragraphs[0]
        p.text = "四、慧企平台月报"
        set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

        txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2), Cm(13), Cm(1))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # # 自动换行
        # tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "本月运维情况分析---集群资源可分配容量"
        set_text_frame_paragraph_format(tfp, font_size=16)

        # cpu资源使用情况
        create_column_chart(slide=slide,
                            categories=cluster_names,
                            series={"已分配": [month_analyse_data[-1][key]["cpu"][0] for key in cluster_names],
                                    "未分配": [1 - month_analyse_data[-1][key]["cpu"][0] for key in cluster_names]},
                            position_and_size=[Cm(1.1), Cm(4.5), Cm(11.5), Cm(8.5)],
                            chart_title="CPU资源使用情况",
                            chart_has_legend=True,
                            column_chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                            show_data_labels=True,
                            y_axis_show_percent=True)
        # 内存资源使用情况
        create_column_chart(slide=slide,
                            categories=cluster_names,
                            series={"已分配": [month_analyse_data[-1][key]["memory"][0] for key in cluster_names],
                                    "未分配": [1 - month_analyse_data[-1][key]["memory"][0] for key in cluster_names]},
                            position_and_size=[Cm(13.6), Cm(4.5), Cm(11.5), Cm(8.5)],
                            chart_title="内存资源使用情况",
                            chart_has_legend=True,
                            column_chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                            show_data_labels=True,
                            y_axis_show_percent=True)

        for cluster_name in cluster_names:
            table_data = cluster_table_data[cluster_name]

            # 定义每个集群（运行情况分析）ppt页数
            # 每页最多存60行数据
            pages = 0
            if len(table_data) % 60:
                pages = len(table_data) // 60 + 1
            else:
                pages = len(table_data) // 60

            for page in range(1, pages + 1):
                # 添加一张幻灯片
                slide = self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

                # 幻灯片标题
                slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
                p = slidePlaceholder.text_frame.paragraphs[0]
                p.text = "四、慧企平台月报"
                set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

                txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2.5), Cm(8), Cm(1))
                # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                tf = txbox.text_frame
                # # 自动换行
                # tf.word_wrap = True
                tfp = tf.paragraphs[0]
                tfp.text = cluster_name + "资源使用情况"
                set_text_frame_paragraph_format(tfp, font_size=16)

                # 资源使用情况折线图
                create_line_chart(slide=slide,
                                  categories=weeks[:data_length],
                                  series={"CPU使用率": cpu_allocated[cluster_name],
                                          "内存使用率": memory_allocated[cluster_name]},
                                  position_and_size=[Cm(0.2), Cm(4), Cm(6), Cm(7.5)],
                                  chart_title="资源使用率",
                                  chart_has_legend=True,
                                  y_axis_show_percent=True)
                if page == 1 and pages == page:
                    if len(table_data) <= 30:
                        # table_1
                        # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_summary[cluster_name]
                        set_text_frame_paragraph_format(tfp)

                    elif len(table_data) > 30:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_summary[cluster_name]
                        set_text_frame_paragraph_format(tfp)
                elif page == 1:
                    # table_1
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.rows[0].heitht = Cm(0.3)
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="F6E7E7")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="CC0000")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])

                    # table_2
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    # 取消表格第一行特殊格式
                    table.first_row = False
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="F6E7E7")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="CC0000")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])
                else:
                    if page == pages:
                        if len(table_data) <= 30:
                            # table_1
                            # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_summary[cluster_name]
                            set_text_frame_paragraph_format(tfp)
                        elif len(table_data) > 30:
                            # table_1
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # table_2
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)
                            # 取消表格第一行特殊格式
                            table.first_row = False
                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_summary[cluster_name]
                            set_text_frame_paragraph_format(tfp)
                    else:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])



        # 添加一张幻灯片
        slide = self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

        # 幻灯片标题
        slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
        p = slidePlaceholder.text_frame.paragraphs[0]
        p.text = "四、慧企平台月报"
        set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

        txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2), Cm(13), Cm(1))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # # 自动换行
        # tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "工作中的问题或困难"
        set_text_frame_paragraph_format(tfp, font_size=16)
