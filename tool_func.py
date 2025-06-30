import pymysql
import json
import yaml
import urllib3
import pptx.table
import pptx.text.text
import pptx.chart.chart
import pptx.presentation
import pptx.slide
import pptx.shapes.base
import pptx.table
import pptx.chart.plot
import pptx.chart.series
import pptx.shapes.graphfrm
import pptx.dml.chtfmt
import pptx.dml.fill
import pptx.dml.line
import pptx.chart.axis
from pptx.util import Pt
import pptx.shapes.placeholder
import pptx.shapes.autoshape
from pptx.dml.color import RGBColor
from kubernetes import client, config
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION, XL_CHART_TYPE, XL_DATA_LABEL_POSITION

urllib3.disable_warnings()


class MySQLHandler:
    def __init__(self, conn):
        if isinstance(conn, pymysql.connections.Connection):
            self._conn = conn
            self._cursor = conn.cursor()
        else:
            raise TypeError

    def select(self, sql):
        self._cursor.execute(sql)
        result = self._cursor.fetchall()
        self._cursor.close()
        self._conn.close()
        return result


def get_deploy_info():
    # 集群context
    context_list = []
    # 集群部署信息
    deploy_info = {}
    # 集群总cpu和内存
    cluster_total_resources = {}
    # 集群中各namespace的资源配额
    cluster_namespaces_quota = {}

    with open("config/kubeconfig/config", encoding="utf-8") as kubeconfig:
        kc = yaml.safe_load(kubeconfig)
        for item in kc.get("contexts"):
            context_list.append(item.get("name"))
    # for context in context_list:
    context = "fcp"
    deploy_info[context] = {}
    cluster_namespaces_quota[context] = {}
    namespaces_names = []
    config.load_kube_config(config_file="config/kubeconfig/config", context=context)
    namespaces = client.CoreV1Api().list_namespace()
    nodes = client.CoreV1Api().list_node()
    cluster_total_resources[context] = {"cpu": 0, "memory": 0}
    # 通过集群各节点资源计算集群worker节点资源总额（总资源减去master资源）
    for node in nodes.items:
        cluster_total_resources[context]["cpu"] += int(node.status.capacity.get("cpu"))
        cluster_total_resources[context]["memory"] += int(node.status.capacity.get("memory").split("Ki")[0])
    if context == "fcp":
        cluster_total_resources[context]["cpu"] -= 48 * 3
        cluster_total_resources[context]["memory"] = cluster_total_resources[context]["memory"] // (
                1024 * 1024) - 376 * 3
    else:
        cluster_total_resources[context]["cpu"] -= 56 * 3
        cluster_total_resources[context]["memory"] = cluster_total_resources[context]["memory"] // (
                1024 * 1024) - 251 * 3
    # 获取集群namespace
    for namespace in namespaces.items:
        if namespace.metadata.name == "cnspnspace-fcp":
            continue
        namespaces_names.append(namespace.metadata.name)
    # 获取集群中各namespace的应用部署信息（通过deploy或statefulset的labels标签获取应用名称）以及资源配额
    for ns in namespaces_names:
        # 集群namespace应用部署信息
        deploy_info[context][ns] = {}
        # 集群namespace资源配额信息
        cluster_namespaces_quota[context][ns] = {"cpu": 0, "memory": 0}
        # 绿区fcp集群的statefulset的部署信息
        if context == "fcp":
            statefulsets = client.AppsV1Api().list_namespaced_stateful_set(namespace=ns)
            for sts in statefulsets.items:
                if sts.metadata.labels.get("app"):
                    deploy_info[context][ns][sts.metadata.labels.get("app")] = sts.spec.replicas
        # 集群deployment的部署信息
        deploys = client.ExtensionsV1beta1Api().list_namespaced_deployment(namespace=ns)
        for deploy in deploys.items:
            # 没label的就舍弃
            if deploy.metadata.labels.get("application_name"):
                deploy_info[context][ns][deploy.metadata.labels.get("application_name")] = deploy.spec.replicas
            elif deploy.metadata.labels.get("app"):
                deploy_info[context][ns][deploy.metadata.labels.get("app")] = deploy.spec.replicas

        # 集群namespace资源配额
        quotas = client.CoreV1Api().list_namespaced_resource_quota(namespace=ns)
        for quota in quotas.items:
            cpu = quota.spec.hard.get("limits.cpu")
            memory = quota.spec.hard.get("limits.memory")
            cluster_namespaces_quota[context][ns]["cpu"] = int(cpu)
            # 单位Gi
            if quota.spec.hard.get("limits.memory")[-2:] == "Gi":
                cluster_namespaces_quota[context][ns]["memory"] = int(memory.split("Gi")[0])
            elif quota.spec.hard.get("limits.memory")[-2:] == "Ti":
                cluster_namespaces_quota[context][ns]["memory"] = int(memory.split("Ti")[0]) * 1024
            elif quota.spec.hard.get("limits.memory")[-2:] == "Mi":
                cluster_namespaces_quota[context][ns]["memory"] = int(memory.split("Mi")[0]) // 1024
            elif quota.spec.hard.get("limits.memory")[-2:] == "Ki":
                cluster_namespaces_quota[context][ns]["memory"] = int(memory.split("Ki")[0]) / 1024 // 1024

    return deploy_info, cluster_total_resources, cluster_namespaces_quota


def get_namespace_project():
    with open('config/config.json', encoding='utf-8') as config:
        config = json.load(config)
    conn = pymysql.connect(**(config["mysql"]))
    mh = MySQLHandler(conn)
    sql = "select PROJECT_ID, PROJECT_NAME, PROJECT_NAME_SPACE from dps_pjm_project"
    return mh.select(sql)


def generate_table_data():
    # 从数据库中扩区项目信息（元组）
    # (('01fd132b9a74487cb3a423d5aa74ff2e', '财务公司头寸管理'), ...)
    namespaces_tuple_info = get_namespace_project()
    # 将数据库的project_id和project_name的元组信息转换为字典
    namespaces_info = {}
    for item in namespaces_tuple_info:
        if item[2]:
            namespaces_info[item[2]] = item[1]
        else:
            namespaces_info[item[0]] = item[1]

    # 从集群获取应用部署信息，集群资源总量，集群namespace资源分配信息

    deploy_info, cluster_total_resources, cluster_namespaces_quota = get_deploy_info()

    # 用于填充运行情况分析table数据
    cluster_table_data = {}
    for key in deploy_info.keys():
        cluster_table_data[key] = [["项目名称", "服务名称", "实例数"]]
        for ns in deploy_info.get(key).keys():
            if ns == ("testtest-fcp" or "cnspnspace-fcp"):
                continue
            postfix = ns.split('-')[-1]
            if postfix == 'fcp' or postfix == 'microservice' or postfix == 'backend':
                for replica_info in deploy_info.get(key).get(ns).items():
                    tmp = []
                    for prj_id in namespaces_info.keys():
                        if prj_id in ns:
                            tmp.append(namespaces_info[prj_id])
                    tmp += list(replica_info)
                    cluster_table_data[key].append(tmp)

        # 合计行
        total = []
        # 统计项目数
        ns_temp = []
        # 统计应用副本数
        replicas_count = 0
        for index in range(1, len(cluster_table_data[key])):
            replicas_count += int(cluster_table_data[key][index][2])
            ns_temp.append(cluster_table_data[key][index][0])
        ns_temp = list(set(ns_temp))
        total.append(str(len(ns_temp)))
        # 统计应用数
        total.append(str(len(cluster_table_data[key]) - 1))
        total.append(str(replicas_count))
        cluster_table_data[key].append(total)

    # 处理有一下table_data，为首列单元格合并做准备
    for key in cluster_table_data.keys():
        tmp_text = ""
        for item in cluster_table_data[key]:
            if item[0] != tmp_text:
                tmp_text = item[0]
            else:
                item[0] = ""

    # 为运行情况分析的扇形图提供数据
    cluster_pie_data = {}
    for cluster in cluster_namespaces_quota.keys():
        # cpu: 核，memory: Gi
        cluster_pie_data[cluster] = {
            "cpu": {"allocated": 0, "total": 0},
            "memory": {"allocated": 0, "total": 0}
        }
        allocated_cpu = 0
        allocated_memory = 0
        for ns in cluster_namespaces_quota.get(cluster).keys():
            allocated_cpu += cluster_namespaces_quota.get(cluster).get(ns).get("cpu")
            allocated_memory += cluster_namespaces_quota.get(cluster).get(ns).get("memory")

        cluster_pie_data[cluster]["cpu"]["allocated"] = allocated_cpu
        cluster_pie_data[cluster]["memory"]["allocated"] = allocated_memory
        cluster_pie_data[cluster]["cpu"]["total"] = cluster_total_resources.get(cluster).get("cpu")
        cluster_pie_data[cluster]["memory"]["total"] = cluster_total_resources.get(cluster).get("memory")

    return cluster_pie_data, cluster_table_data


# 调整末尾新增的幻灯片到指定位置，针对《运行情况分析》
def insert_slide(insert_index: int, insert_length: int, prs: pptx.presentation.Presentation):
    xml_slides = prs.slides._sldIdLst
    slideIdList = list(xml_slides)
    for i in range(insert_length):
        item = slideIdList.pop()
        xml_slides.insert(insert_index, item)


# 先写内容，后改样式，否则可能出现样式不生效问题，如字体之类的？？？
def set_cell_format(cell: pptx.table._Cell,
                    cell_background_color=None,
                    margin_left=None,
                    margin_right=None,
                    margin_top=None,
                    margin_bottom=None,
                    vertical_anchor=MSO_ANCHOR.MIDDLE):
    # cell中文字相对cell边框的距离，类似与html中的padding
    cell.margin_left = margin_left
    cell.margin_right = margin_right
    cell.margin_top = margin_top
    cell.margin_bottom = margin_bottom

    # 垂直对齐
    cell.vertical_anchor = vertical_anchor

    # 颜色填充
    if cell_background_color:
        # 填充cell前景色
        # cell.fill.patterned()  # 图案填充（可以设置cell前景色和背景色）
        cell.fill.solid()  # 纯色填充（只能设置cell前景色，但是对cell中的text而言也是背景）
        # cell.fill.back_color.rgb = RGBColor.from_string(cell_background_color)
        cell.fill.fore_color.rgb = RGBColor.from_string(cell_background_color)
    # else:
    #     # cell颜色无填充
    #     cell.fill.background()


# 设置text_frame的paragraph文字格式
def set_text_frame_paragraph_format(
        tfp: pptx.text.text._Paragraph,
        font_name="Microsoft YaHei",
        font_size=8,
        font_bold=False,
        align=PP_ALIGN.LEFT,
        font_color="000000"):
    # 语言设置，NONE表示移除所有语言设置
    tfp.font.language_id = MSO_LANGUAGE_ID.NONE
    # 字体
    tfp.font.name = font_name
    # 字体大小
    tfp.font.size = Pt(int(font_size))
    # 是否加粗
    tfp.font.bold = font_bold
    # 水平对齐
    tfp.alignment = align

    # 用RGB表示字体颜色（两种方式）
    # tfp.font.color.rgb = RGBColor.from_string(fontColor)
    # 设置前景色或背景色之前需要执行
    # tfp.font.fill.patterned()  # 图案填充（可以设置前景色和背景色）
    tfp.font.fill.solid()  # 纯色填充（只能设置前景色）
    # 字体前景色(就是字体颜色）
    tfp.font.fill.fore_color.rgb = RGBColor.from_string(font_color)
    # 字体背景色（就是文字本身的背景，不是整个cell）
    # tfp.font.fill.back_color.rgb = RGBColor.from_string(cell_background_color)
    # 字体颜色透明度
    # tfp.font.color.brightness = -1  # 取值范围-1~1，暗->亮


# 扇形图（周报）
def create_pie_chart(slide,
                     categories: list,
                     series: dict,
                     position_and_size: list,
                     chart_title: str,
                     chart_title_font_size=10,
                     chart_title_font_name="Microsoft YaHei",
                     chart_text_font_size=10,
                     chart_text_font_name="Microsoft YaHei",
                     chart_has_legend=False,
                     chart_legend_position=XL_LEGEND_POSITION.BOTTOM,
                     pie_chart_type=XL_CHART_TYPE.PIE,
                     data_labels_position=XL_DATA_LABEL_POSITION.CENTER,
                     data_labels_font_size=8,
                     data_labels_font_name="Microsoft YaHei"):
    """
    :param slide: 幻灯片实例
    :param categories: 分类
    :param series: 数据系列/分类
    :param position_and_size: 图形的位置和大小
            for example:
              position_and_size = [ Cm(1), Cm(2) , Cm(3), Cm(4)]
            means:
             position:
               left = Cm(1)
               top = Cm(2)
             size:
               width = Cm(3)
               height = Cm(4)
    :param chart_title: chart的标题内容（顶部居中）
    :param chart_title_font_size: 标题字体大小
    :param chart_title_font_name: 标题字体名称
    :param chart_text_font_size: 除标题外的其他字体大小
    :param chart_text_font_name: 除标题外的其他字体名称
    :param chart_has_legend: 是否显示图例
    :param chart_legend_position: 图例位置
    :param pie_chart_type: pie chart 类型
    :param data_labels_position: data_labels位置
    :param data_labels_font_size: 数据标签字体大小
    :param data_labels_font_name: 数据标签字体名称
    :return:
    """
    # cpu 扇形图
    chart_data = ChartData()
    chart_data.categories = categories
    # 可以设置多组数据填入扇形图编辑时的excel中，但只有第一组数据会显示
    for key in series.keys():
        chart_data.add_series(key, tuple(series[key]))
    gf = slide.shapes.add_chart(pie_chart_type, *position_and_size, chart_data)
    chart = gf.chart

    # 设置图例说明(扇形图的图例是categories)
    if chart_has_legend:
        chart.has_legend = chart_has_legend
        chart.legend.position = chart_legend_position
        chart.legend.include_in_layout = False

    chart.font.name = chart_text_font_name
    chart.font.size = Pt(chart_text_font_size)

    # 设置标题（不设置，默认值是add_series中的列标题"cpu"）
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = chart_title
    set_text_frame_paragraph_format(chart.chart_title.text_frame.paragraphs[0],
                                    font_size=chart_title_font_size,
                                    font_name=chart_title_font_name)

    chart.plots[0].has_data_labels = True
    chart_data_labels = chart.plots[0].data_labels
    chart_data_labels.number_format = '0%'
    chart_data_labels.position = data_labels_position
    chart_data_labels.font.size = Pt(data_labels_font_size)
    chart_data_labels.font.name = data_labels_font_name


# 柱状图（月报）
def create_column_chart(slide,
                        categories: list,
                        series: dict,
                        position_and_size: list,
                        chart_title: str,
                        chart_title_font_size=14,
                        chart_title_font_name="Microsoft YaHei",
                        chart_xy_axis_text_font_size=8,
                        chart_xy_axis_text_font_name="Microsoft YaHei",
                        chart_has_legend=False,
                        chart_legend_position=XL_LEGEND_POSITION.BOTTOM,
                        column_chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
                        column_color="ECBBAF",
                        show_data_labels=False,
                        data_labels_position=XL_DATA_LABEL_POSITION.CENTER,
                        data_labels_font_size=6,
                        data_labels_font_name="Microsoft YaHei",
                        y_axis_show_percent=False):
    """
    :param slide: 幻灯片实例
    :param categories: 分类（横坐标）
    :param series: 数据系列/分类
    :param position_and_size: 图形的位置和大小
            for example:
              position_and_size = [ Cm(1), Cm(2) , Cm(3), Cm(4)]
            means:
             position:
               left = Cm(1)
               top = Cm(2)
             size:
               width = Cm(3)
               height = Cm(4)
    :param chart_title: chart的标题内容（顶部居中）
    :param chart_title_font_size: 标题字体大小
    :param chart_title_font_name: 标题字体名称
    :param chart_xy_axis_text_font_size: 横纵坐标字体大小
    :param chart_xy_axis_text_font_name: 横纵坐标字体名称
    :param chart_has_legend: 是否显示图例
    :param chart_legend_position: 图例(系列)位置(need: chart_has_legend = True)
    :param column_chart_type: 柱状图类型
           for example:
             column_chart_type = XL_CHART_TYPE.COLUMN_STACKED_100
    :param column_color: 柱状图颜色
    :param show_data_labels: 是否在图形上显示数据
    :param data_labels_position: data_labels位置(need: show_data_labels = True)
    :param data_labels_font_size: 数据标签字体大小(need:show_data_labels = True)
    :param data_labels_font_name: 数据标签字体名称(need:show_data_labels = True)
    :param y_axis_show_percent: Y轴是否百分比显示
    :return:
    """
    # 初始化柱状图数据
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for key in series.keys():
        chart_data.add_series(key, tuple(series[key]))
    # 画柱状图
    gf = slide.shapes.add_chart(
        column_chart_type,
        *position_and_size,
        chart_data
    )  # type: pptx.shapes.graphfrm.GraphicFrame
    chart = gf.chart  # type: pptx.chart.chart.Chart

    # 设置图例说明（会在图中标识已分配、未分配的颜色说明）
    if chart_has_legend:
        chart.has_legend = chart_has_legend
        chart.legend.position = chart_legend_position
        chart.legend.include_in_layout = False

    # 设置标题
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = chart_title
    set_text_frame_paragraph_format(chart.chart_title.text_frame.paragraphs[0],
                                    font_size=chart_title_font_size,
                                    font_name=chart_title_font_name)
    # 设置横纵坐标的字体
    chart.font.name = chart_xy_axis_text_font_name
    chart.font.size = Pt(chart_xy_axis_text_font_size)

    plot = chart.plots[0]  # type: pptx.chart.plot._BasePlot

    # 是否显示data_labels
    if show_data_labels:
        plot.has_data_labels = True
        chart_data_labels = plot.data_labels
        chart_data_labels.font.name = data_labels_font_name
        chart_data_labels.font.size = Pt(data_labels_font_size)
        chart_data_labels.position = data_labels_position
        if y_axis_show_percent:
            chart_data_labels.number_format = '0%'

    # 获取分类的数据（系列1、系列2、……）(此处我们为已分配和未分配)
    # 已分配
    sri = plot.series[0]  # type: pptx.chart.series.BarSeries
    cf = sri.format  # type: pptx.dml.chtfmt.ChartFormat

    # 柱形轮廓线格式
    lf = cf.line  # type: pptx.dml.line.LineFormat
    lf.color.rgb = RGBColor.from_string("7C7D7C")
    lf.width = 3600

    # 已分配->柱形填充
    ff = cf.fill  # type: pptx.dml.fill.FillFormat
    ff.solid()
    ff.fore_color.rgb = RGBColor.from_string(column_color)

    # 未分配
    if len(series.keys()) == 2 and column_chart_type == XL_CHART_TYPE.COLUMN_STACKED_100 or column_chart_type == XL_CHART_TYPE.COLUMN_STACKED:
        sri = plot.series[1]  # type: pptx.chart.series.BarSeries
        cf = sri.format  # type: pptx.dml.chtfmt.ChartFormat

        # 柱形轮廓线格式
        lf = cf.line  # type: pptx.dml.line.LineFormat
        lf.color.rgb = RGBColor.from_string("7C7D7C")
        lf.width = 3600

        # 未分配->柱形填充
        ff = cf.fill  # type: pptx.dml.fill.FillFormat
        ff.solid()
        ff.fore_color.rgb = RGBColor.from_string("FFFFFF")

    # 轴
    x_axis = chart.category_axis  # type: pptx.chart.axis.CategoryAxis
    y_axis = chart.value_axis  # type: pptx.chart.axis.ValueAxis

    if y_axis_show_percent:
        # 设置Y轴数据区间
        y_axis.minimum_scale = 0
        y_axis.maximum_scale = 1

        # 设置Y轴数据格式为百分比
        y_axis.tick_labels.number_format = "0%"

    # 设置Y轴的网格线(可以设置主要和次要网格线)
    mg = y_axis.major_gridlines  # type: pptx.chart.axis.MajorGridlines
    lf = mg.format.line  # type: pptx.dml.line.LineFormat
    lf.color.rgb = RGBColor.from_string("D9D9D9")
    lf.width = 3600

    # 取消轴的线条填充
    x_axis.format.line.fill.background()
    y_axis.format.line.fill.background()


# 曲线图（月报）
def create_line_chart(slide,
                      categories: list,
                      series: dict,
                      position_and_size: list,
                      chart_title: str,
                      chart_title_font_size=14,
                      chart_title_font_name="Microsoft YaHei",
                      chart_xy_axis_text_font_size=8,
                      chart_xy_axis_text_font_name="Microsoft YaHei",
                      chart_has_legend=False,
                      chart_legend_position=XL_LEGEND_POSITION.BOTTOM,
                      line_chart_type=XL_CHART_TYPE.LINE,
                      show_data_labels=False,
                      data_labels_position=XL_DATA_LABEL_POSITION.CENTER,
                      data_labels_font_size=6,
                      data_labels_font_name="Microsoft YaHei",
                      y_axis_show_percent=False):
    """
    :param slide: 幻灯片实例
    :param categories: 分类（横坐标）
    :param series: 数据系列/分类
    :param position_and_size: 图形的位置和大小
            for example:
              position_and_size = [ Cm(1), Cm(2) , Cm(3), Cm(4)]
            means:
             position:
               left = Cm(1)
               top = Cm(2)
             size:
               width = Cm(3)
               height = Cm(4)
    :param chart_title: chart的标题内容（顶部居中）
    :param chart_title_font_size: 标题字体大小
    :param chart_title_font_name: 标题字体名称
    :param chart_xy_axis_text_font_size: 横纵坐标字体大小
    :param chart_xy_axis_text_font_name: 横纵坐标字体名称
    :param chart_has_legend: 是否显示图例
    :param chart_legend_position: 图例(系列)位置(need: chart_has_legend = True)
    :param line_chart_type: 曲线图类型
           for example:
             line_chart_type = XL_CHART_TYPE.LINE
    :param show_data_labels: 是否在图形上显示数据
    :param data_labels_position: data_labels位置(need: show_data_labels = True)
    :param data_labels_font_size: 数据标签字体大小(need:show_data_labels = True)
    :param data_labels_font_name: 数据标签字体名称(need:show_data_labels = True)
    :param y_axis_show_percent: Y轴是否百分比显示
    :return:
    """

    # 资源使用情况折线图
    chart_data = ChartData()
    chart_data.categories = categories

    for key in series.keys():
        chart_data.add_series(key, tuple(series[key]))
    gf = slide.shapes.add_chart(
        line_chart_type,
        *position_and_size,
        chart_data
    )  # type: pptx.shapes.graphfrm.GraphicFrame
    chart = gf.chart  # type: pptx.chart.chart.Chart

    plot = chart.plots[0]  # type: pptx.chart.plot._BasePlot

    if show_data_labels:
        plot.has_data_labels = True
        chart_data_labels = plot.data_labels
        chart_data_labels.font.name = data_labels_font_name
        chart_data_labels.font.size = Pt(data_labels_font_size)
        chart_data_labels.position = data_labels_position
        if y_axis_show_percent:
            chart_data_labels.number_format = '0%'

    # 设置图例说明（会在图中标识已分配、未分配的颜色说明）
    chart.has_legend = True
    chart.legend.position = chart_legend_position
    chart.legend.include_in_layout = False

    # 横纵坐标字体
    chart.font.name = chart_xy_axis_text_font_name
    chart.font.size = Pt(chart_xy_axis_text_font_size)

    # 设置标题（不设置，默认值是add_series中的列标题"cpu"）-> "ChartTitle" has no attribute "width"??
    chart.has_title = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = chart_title
    set_text_frame_paragraph_format(chart.chart_title.text_frame.paragraphs[0],
                                    font_size=chart_title_font_size,
                                    font_name=chart_title_font_name)

    # 轴
    x_axis = chart.category_axis  # type: pptx.chart.axis.CategoryAxis
    y_axis = chart.value_axis  # type: pptx.chart.axis.ValueAxis

    if y_axis_show_percent:
        # 设置Y轴数据区间
        y_axis.minimum_scale = 0
        y_axis.maximum_scale = 1

        # 设置Y轴数据格式为百分比
        y_axis.tick_labels.number_format = "0%"

    # 设置Y轴的网格线(可以设置主要和次要网格线)
    mg = y_axis.major_gridlines  # type: pptx.chart.axis.MajorGridlines
    lf = mg.format.line  # type: pptx.dml.line.LineFormat
    lf.color.rgb = RGBColor.from_string("D9D9D9")
    lf.width = 3600

    # 取消轴的线条填充
    x_axis.format.line.fill.background()
    y_axis.format.line.fill.background()
