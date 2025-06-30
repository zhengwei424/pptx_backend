import pptx
import pptx.slide
import pptx.table
import pptx.text.text
import pptx.shapes.base
import pptx.chart.chart
from pptx.util import Cm
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from tool_func import set_text_frame_paragraph_format, set_cell_format, create_pie_chart


# presentation -> slide -> shapes -> placeholder,graphfrm -> chart(table -> cell) -> text_frame -> paragraphs -> font
#            |      +-> placeholder
#            |->slide_master -> slide_layout
# https://mhammond.github.io/pywin32/html/com/win32com/HTML/QuickStartClientCom.html


class SetWeaklyReport(object):
    def __init__(self, prs):
        self._prs = prs
        self._current_slide_idx = 0  # 当前操作的幻灯片索引

    # 将新添加的slide（默认追加在slides末尾）插入到指定索引之前
    def insert_slide(self, insert_index: int):
        xml_slides = self._prs.slides._sldIdLst
        slideIdList = list(xml_slides)
        item = slideIdList.pop()
        xml_slides.insert(insert_index, item)

    # 1. 运维工作统计（次数）
    def slide_1(self, events_count: list):
        """运维工作统计"""
        self._current_slide_idx = 8
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(3, 1),
                    tb.cell(3, 2),
                    tb.cell(3, 3),
                    tb.cell(3, 4),
                    tb.cell(3, 5),
                    tb.cell(3, 6),
                ]
                for index in range(len(cells)):
                    cells[index].text = str(events_count[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=11,
                                                    align=PP_ALIGN.CENTER
                                                    )

    # 2. 巡检
    def slide_2(self, inspect_data: list):
        """巡检"""
        if not inspect_data:
            inspect_data = ["11", "0", "6", "√", "√", "√", "√", "√", "√", "√", "√", "√", "√", "", "", "√", ""]
        self._current_slide_idx = 9
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(3, 1),
                    tb.cell(3, 2),
                    tb.cell(3, 3),
                    tb.cell(3, 5),
                    tb.cell(4, 5),
                    tb.cell(3, 7),
                    tb.cell(4, 7),
                    tb.cell(3, 9),
                    tb.cell(4, 9),
                    tb.cell(3, 11),
                    tb.cell(4, 11),
                    tb.cell(3, 13),
                    tb.cell(4, 13),
                    tb.cell(3, 15),
                    tb.cell(4, 15),
                    tb.cell(3, 17),
                ]
                for index in range(len(cells)):
                    cells[index].text = str(inspect_data[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=11,
                                                    align=PP_ALIGN.CENTER
                                                    )

    # 3. 变更
    def slide_3(self, change_data: list):
        """变更"""
        self._current_slide_idx = 10
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        # 行数初始化
        row = 5
        if len(change_data) + 1 > row:
            row = len(change_data) + 1
        # 创建table
        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
        x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
        table = slide.shapes.add_table(row, 7, x, y, cx, cy).table
        table.columns[0].width = Cm(2)
        table.columns[1].width = Cm(3.5)
        table.columns[2].width = Cm(3.5)
        table.columns[3].width = Cm(3.5)
        table.columns[4].width = Cm(3.5)
        table.columns[5].width = Cm(3.5)
        table.columns[6].width = Cm(3.5)
        first_row_cells = [
            table.cell(0, 0),
            table.cell(0, 1),
            table.cell(0, 2),
            table.cell(0, 3),
            table.cell(0, 4),
            table.cell(0, 5),
            table.cell(0, 6)
        ]
        table_header = [
            "专业",
            "变更类别",
            "变更内容",
            "变更影响",
            "变更时间",
            "资源支持",
            "目前进展"
        ]
        for index in range(len(first_row_cells)):
            p = first_row_cells[index].text_frame.paragraphs[0]
            p.text = table_header[index]
            set_text_frame_paragraph_format(p,
                                            font_size=16,
                                            font_bold=True,
                                            align=PP_ALIGN.CENTER,
                                            font_color="FFFFFF")
        if change_data:
            for row_idx in range(1, len(table.rows)):
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    if row_idx <= len(change_data):
                        tfp.text = change_data[row_idx - 1][col_idx]
                        set_text_frame_paragraph_format(tfp,
                                                        font_size=11,
                                                        align=PP_ALIGN.LEFT)
                    else:
                        break

    # 4. 支撑发版
    def slide_4(self, release_data: list):
        """支撑发版"""
        # 行数
        row = 15
        # 表格页数初始化
        pages = len(release_data) // (row - 1) + 1
        # if len(release_data) == 0:
        #     pages = 1
        # elif len(release_data) % (row - 1):
        #     pages = len(release_data) // (row - 1) + 1
        # else:
        #     pages = len(release_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2.4 支撑发版"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
            table.columns[0].width = Cm(3)
            table.columns[1].width = Cm(4)
            table.columns[2].width = Cm(3)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(5)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
            ]
            table_header = [
                "专业",
                "发版时间",
                "次数",
                "工作内容",
                "异常情况处理",
            ]
            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not release_data:
                    break
                data = release_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 5. 资源权限管理
    def slide_5(self, permission_management_data: list):
        """资源权限管理"""
        # 行数
        row = 9
        # 表格页数初始化
        pages = len(permission_management_data) // (row - 1) + 1
        # if len(permission_management_data) == 0:
        #     pages = 1
        # elif len(permission_management_data) % (row - 1):
        #     pages = len(permission_management_data) // (row - 1) + 1
        # else:
        #     pages = len(permission_management_data) // (row - 1)
        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2.5 资源权限管理"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(3)
            table.columns[3].width = Cm(4)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]
            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not permission_management_data:
                    break
                data = permission_management_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 6. 配合操作及排障
    def slide_6(self, cooperation_data: list):
        """配合操作及排障"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = len(cooperation_data) // (row - 1) + 1
        # if len(cooperation_data) == 0:
        #     pages = 1
        # elif len(cooperation_data) % (row - 1):
        #     pages = len(cooperation_data) // (row - 1) + 1
        # else:
        #     pages = len(cooperation_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2.6 配合操作及排障"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)
            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(3)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]
            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not cooperation_data:
                    break
                data = cooperation_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 7. 问题及告警
    def slide_7(self, problem_data: list):
        """问题及告警"""
        self._current_slide_idx += 1
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        # 行数
        row = 5
        if len(problem_data) + 1 > row:
            row = len(problem_data) + 1

        # 创建table
        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
        x, y, cx, cy = Cm(0.5), Cm(8), Cm(1), Cm(0.5)
        table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
        table.columns[0].width = Cm(2)
        table.columns[1].width = Cm(5)
        table.columns[2].width = Cm(5)
        table.columns[3].width = Cm(5)
        table.columns[4].width = Cm(5)
        first_row_cells = [
            table.cell(0, 0),
            table.cell(0, 1),
            table.cell(0, 2),
            table.cell(0, 3),
            table.cell(0, 4),
        ]
        table_header = [
            "专业",
            "问题描述",
            "处理结果",
            "原因分析",
            "后续建议",
        ]
        for index in range(len(first_row_cells)):
            p = first_row_cells[index].text_frame.paragraphs[0]
            p.text = table_header[index]
            set_text_frame_paragraph_format(p,
                                            font_size=16,
                                            font_bold=True,
                                            align=PP_ALIGN.CENTER,
                                            font_color="FFFFFF")

        if problem_data:
            for row_idx in range(1, len(table.rows)):
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    if row_idx <= len(problem_data):
                        tfp.text = problem_data[row_idx - 1][col_idx]
                        set_text_frame_paragraph_format(tfp,
                                                        font_size=11,
                                                        align=PP_ALIGN.LEFT)
                    else:
                        break

    # 8. 运行情况分析
    def slide_8(self, cluster_pie_data, cluster_table_data):
        """运行情况分析"""
        self._current_slide_idx += 3
        # 传入集群应用部署信息
        for cluster in cluster_table_data.keys():
            table_data = cluster_table_data[cluster]

            # 定义每个集群（运行情况分析）ppt页数
            # 每页最多存60行数据
            pages = 0
            if len(table_data) % 60:
                pages = len(table_data) // 60 + 1
            else:
                pages = len(table_data) // 60

            # 集群名称
            cluster_name = ""
            if cluster == "fcp":
                cluster_name = "FCP业务集群"
            # elif cluster == "fcp-inner-microservice":
            #     cluster_name = "内网微服务区"
            # elif cluster == "fcp-inner-backend":
            #     cluster_name = "内网后台区"
            # elif cluster == "fcp-outer-microservice":
            #     cluster_name = "外网微服务区"
            # elif cluster == "fcp-outer-backend":
            #     cluster_name = "外网后台区"
            # 集群合计信息
            total_count = table_data[-1]
            # 集群资源已分配数值
            cluster_allocated_cpu = cluster_pie_data.get(cluster).get("cpu").get("allocated")
            cluster_allocated_memory = cluster_pie_data.get(cluster).get("memory").get("allocated")
            # 集群资源总数值
            cluster_total_cpu = cluster_pie_data.get(cluster).get("cpu").get("total")
            cluster_total_memory = cluster_pie_data.get(cluster).get("memory").get("total")
            # 集群资源已分配比率
            cpu_allocated_rate = float("%.2f" % (cluster_allocated_cpu / cluster_total_cpu))
            memory_allocated_rate = float("%.2f" % (cluster_allocated_memory / cluster_total_memory))

            for page in range(1, pages + 1):
                # 添加一张幻灯片
                self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
                self.insert_slide(self._current_slide_idx + 1)
                self._current_slide_idx += 1

                slide = self._prs.slides[self._current_slide_idx]
                # 幻灯片标题
                slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
                p = slidePlaceholder.text_frame.paragraphs[0]
                p.text = "2.11慧企运行情况分析–集群资源使用情况"
                set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

                # cpu 扇形图
                create_pie_chart(slide=slide,
                                 categories=["已分配", "未分配"],
                                 series={"cpu": [cpu_allocated_rate, 1 - cpu_allocated_rate]},
                                 position_and_size=[Cm(0.1), Cm(2.2), Cm(6), Cm(6)],
                                 chart_title=cluster_name + " CPU(%)",
                                 chart_has_legend=True)

                # memory 扇形图
                create_pie_chart(slide=slide,
                                 categories=["已分配", "未分配"],
                                 series={"memory": [memory_allocated_rate, 1 - memory_allocated_rate]},
                                 position_and_size=[Cm(0.1), Cm(8.3), Cm(6), Cm(6)],
                                 chart_title=cluster_name + " 内存(%)",
                                 chart_has_legend=True)

                # 添加table
                if page == 1 and pages == page:
                    if len(table_data) <= 30:
                        # table_1
                        # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="E8EFF3")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="3494BA")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_name + "部署" + total_count[0] + "个应用项目，包含" + \
                                   total_count[1] + "个服务，运行实例" + total_count[2] + \
                                   "个。CPU已分配" + str(cluster_allocated_cpu) + "/" + str(
                            cluster_total_cpu) + "，占比" + \
                                   "{:.0%}".format(cpu_allocated_rate) + "，内存已分配" + \
                                   str(cluster_allocated_memory) + "/" + str(cluster_total_memory) + \
                                   "，占比" + "{:.0%}".format(memory_allocated_rate)
                        set_text_frame_paragraph_format(tfp)

                    elif len(table_data) > 30:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="E8EFF3")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="3494BA")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="E8EFF3")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="3494BA")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_name + "部署" + total_count[0] + "个应用项目，包含" + \
                                   total_count[1] + "个服务，运行实例" + total_count[2] + \
                                   "个。CPU已分配" + str(cluster_allocated_cpu) + "/" + str(
                            cluster_total_cpu) + "，占比" + \
                                   "{:.0%}".format(cpu_allocated_rate) + "，内存已分配" + \
                                   str(cluster_allocated_memory) + "/" + str(cluster_total_memory) + \
                                   "，占比" + "{:.0%}".format(memory_allocated_rate)
                        set_text_frame_paragraph_format(tfp)
                elif page == 1:
                    # table_1
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.rows[0].heitht = Cm(0.3)
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="E8EFF3")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="3494BA")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])

                    # table_2
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    # 取消表格第一行特殊格式
                    table.first_row = False
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="E8EFF3")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="3494BA")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])

                else:
                    if page == pages:
                        if len(table_data) <= 30:
                            # table_1
                            # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="E8EFF3")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="3494BA")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_name + "部署" + total_count[0] + "个应用项目，包含" + \
                                       total_count[1] + "个服务，运行实例" + total_count[2] + \
                                       "个。CPU已分配" + str(cluster_allocated_cpu) + "/" + str(
                                cluster_total_cpu) + "，占比" + \
                                       "{:.0%}".format(cpu_allocated_rate) + "，内存已分配" + \
                                       str(cluster_allocated_memory) + "/" + str(cluster_total_memory) + \
                                       "，占比" + "{:.0%}".format(memory_allocated_rate)
                            set_text_frame_paragraph_format(tfp)
                        elif len(table_data) > 30:
                            # table_1
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="E8EFF3")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="3494BA")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # table_2
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)
                            # 取消表格第一行特殊格式
                            table.first_row = False
                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="E8EFF3")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="3494BA")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_name + "部署" + total_count[0] + "个应用项目，包含" + \
                                       total_count[1] + "个服务，运行实例" + total_count[2] + \
                                       "个。CPU已分配" + str(cluster_allocated_cpu) + "/" + str(
                                cluster_total_cpu) + "，占比" + \
                                       "{:.0%}".format(cpu_allocated_rate) + "，内存已分配" + \
                                       str(cluster_allocated_memory) + "/" + str(cluster_total_memory) + \
                                       "，占比" + "{:.0%}".format(memory_allocated_rate)
                            set_text_frame_paragraph_format(tfp)
                    else:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="E8EFF3")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="3494BA")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="E8EFF3")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="3494BA")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

    # 9. 下周工作计划
    def slide_9(self, working_plan_data: list):
        self._current_slide_idx += 4
        slide = self._prs.slides[self._current_slide_idx]  # type: pptx.slide.Slide
        # 行数
        row = 6
        if len(working_plan_data) + 1 > row:
            row = len(working_plan_data) + 1

        # 创建table
        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
        x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
        table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
        table.columns[0].width = Cm(2)
        table.columns[1].width = Cm(3)
        table.columns[2].width = Cm(5)
        table.columns[3].width = Cm(5)
        table.columns[4].width = Cm(5)
        table.columns[5].width = Cm(3)
        first_row_cells = [
            table.cell(0, 0),
            table.cell(0, 1),
            table.cell(0, 2),
            table.cell(0, 3),
            table.cell(0, 4),
            table.cell(0, 5),
        ]
        table_header = [
            "序号",
            "所属专业",
            "工作类别",
            "工作内容",
            "工作进度",
            "后续安排"
        ]
        for index in range(len(first_row_cells)):
            p = first_row_cells[index].text_frame.paragraphs[0]
            p.text = table_header[index]
            set_text_frame_paragraph_format(p,
                                            font_size=16,
                                            font_bold=True,
                                            align=PP_ALIGN.CENTER,
                                            font_color="FFFFFF")

        if working_plan_data:
            for row_idx in range(1, len(table.rows)):
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    if row_idx <= len(working_plan_data):
                        tfp.text = working_plan_data[row_idx - 1][col_idx]
                        set_text_frame_paragraph_format(tfp,
                                                        font_size=11,
                                                        align=PP_ALIGN.LEFT)
                    else:
                        break
