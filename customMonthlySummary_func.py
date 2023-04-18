import pptx.presentation
import pptx.slide
import pptx.shapes.base
import pptx.table
import pptx.chart.chart
import pptx.chart.plot
import pptx.chart.series
import pptx.shapes.graphfrm
import pptx.dml.chtfmt
import pptx.dml.fill
import pptx.dml.line
import pptx.chart.axis
import pptx.shapes.placeholder
import pptx.shapes.autoshape
from tool_func import set_text_frame_paragraph_format, set_cell_format, create_line_chart, create_column_chart
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm


# presentation -> slide -> shapes -> placeholder,graphfrm -> chart(table -> cell) -> text_frame -> paragraphs -> font
#            |      +-> placeholder                            | - > plot(感觉是扇形的扇子，柱形的柱子) -> categories
#            |                                                          | -> series
#            |->slide_master -> slide_layout
# 运维工作统计


class SetMonthlySummaryReport(object):
    def __init__(self, prs):
        self._prs = prs  # type: pptx.presentation.Presentation
        self._current_slide_idx = 0  # 当前操作的幻灯片索引

    def insert_slide(self, insert_index: int):
        xml_slides = self._prs.slides._sldIdLst
        slideIdList = list(xml_slides)
        item = slideIdList.pop()
        xml_slides.insert(insert_index, item)

    # 巡检
    def slide_1(self, monthly_summary_inspect_data: list):
        """巡检"""
        self._current_slide_idx = 5
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(1, 1),  # 巡检次数
                    tb.cell(1, 2),  # 提交报告次数
                    tb.cell(1, 3),  # 正常次数
                    tb.cell(1, 4),  # 异常次数
                ]
                for index in range(len(cells)):
                    cells[index].text = str(monthly_summary_inspect_data[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=12,
                                                    align=PP_ALIGN.CENTER
                                                    )

    # 运维工作统计
    def slide_2(self, monthly_summary_event_count: list, work_summary: list):
        """运维工作统计"""
        self._current_slide_idx = 6
        slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide
        shape = slide.shapes[0]  # type: pptx.shapes.base.BaseShape
        # 填充表格数据
        for shape in slide.shapes:
            # 判断shape是否是表格（即找到需要修改的表格）
            if shape.has_table:
                tb = shape.table  # type: pptx.table.Table
                cells = [
                    tb.cell(1, 1),  # 变更
                    tb.cell(1, 2),  # 资源及权限管理
                    tb.cell(1, 3),  # 配合操作
                    tb.cell(1, 4),  # 支撑发版
                    tb.cell(1, 5),  # 故障及问题处理
                ]
                for index in range(len(cells)):
                    cells[index].text = str(monthly_summary_event_count[index])
                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                    font_size=10,
                                                    align=PP_ALIGN.CENTER
                                                    )

        # 初始化柱状图数据
        create_column_chart(slide=slide,
                            categories=['变更', '资源权限管理', '配合操作', '支撑发版', '故障和问题处理'],
                            series={"运维工作统计": monthly_summary_event_count},
                            position_and_size=[Cm(1.1), Cm(4.5), Cm(10), Cm(8.5)],
                            chart_title="运维工作统计")

        # 总结本月主要工作
        text_list = work_summary
        txbox = slide.shapes.add_textbox(Cm(13.5), Cm(5.5), Cm(11), Cm(5))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # 自动换行
        tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "\n●".join(text_list)
        set_text_frame_paragraph_format(tfp, font_size=10, font_bold=True, font_color="FF0000")

    # 变更
    def slide_3(self, monthly_summary_change_data):
        """变更"""
        # 表格行数
        row = 12
        # 表格页数初始化
        pages = 0
        if len(monthly_summary_change_data) == 0:
            pages = 1
        elif len(monthly_summary_change_data) % (row - 1):
            pages = len(monthly_summary_change_data) // (row - 1) + 1
        else:
            pages = len(monthly_summary_change_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-变更"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)
            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 2, x, y, cx, cy).table
            table.columns[0].width = Cm(1.3)
            table.columns[1].width = Cm(22)
            # 取消第一行特殊格式
            table.first_row = False
            # 合并第一列
            table.cell(0, 0).merge(table.cell(11, 0))

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            p = table.cell(0, 0).text_frame.paragraphs[0]
            p.text = "联通云变更"
            set_text_frame_paragraph_format(p,
                                            font_size=12,
                                            font_bold=False,
                                            align=PP_ALIGN.CENTER,
                                            font_color="000000")
            set_cell_format(table.cell(0, 0),
                            cell_background_color="F6E7E7")
            for row_idx in range(0, len(table.rows)):
                if not monthly_summary_change_data:
                    break
                data = monthly_summary_change_data.pop(0)
                tf = table.cell(row_idx, 1).text_frame
                tfp = table.cell(row_idx, 1).text_frame.paragraphs[0]
                # 自动换行
                tf.word_wrap = True
                tfp.text = data
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

    # 资源权限管理(有序号)
    def slide_4(self, monthly_summary_permission_management_data: list):
        """资源权限管理"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(monthly_summary_permission_management_data) == 0:
            pages = 1
        elif len(monthly_summary_permission_management_data) % (row - 1):
            pages = len(monthly_summary_permission_management_data) // (row - 1) + 1
        else:
            pages = len(monthly_summary_permission_management_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-资源权限管理"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(3)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not monthly_summary_permission_management_data:
                    break
                data = monthly_summary_permission_management_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 配合操作(有序号)
    def slide_5(self, monthly_summary_cooperation_data: list):
        """配合操作"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(monthly_summary_cooperation_data) == 0:
            pages = 1
        elif len(monthly_summary_cooperation_data) % (row - 1):
            pages = len(monthly_summary_cooperation_data) // (row - 1) + 1
        else:
            pages = len(monthly_summary_cooperation_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-配合操作及排障"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)
            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 6, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(3)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(3)
            table.columns[4].width = Cm(8)
            table.columns[5].width = Cm(3)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
                table.cell(0, 5),
            ]
            table_header = [
                "序号",
                "所属专业",
                "环境",
                "需求方",
                "申请内容",
                "完成进度"
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not monthly_summary_cooperation_data:
                    break
                data = monthly_summary_cooperation_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 支撑发版
    def slide_6(self, monthly_summary_release_data: list):
        """支撑发版"""
        # 表格行数
        row = 15
        # 表格页数初始化
        pages = 0
        if len(monthly_summary_release_data) == 0:
            pages = 1
        elif len(monthly_summary_release_data) % (row - 1):
            pages = len(monthly_summary_release_data) // (row - 1) + 1
        else:
            pages = len(monthly_summary_release_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-支撑发版"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
            table.columns[0].width = Cm(3)
            table.columns[1].width = Cm(4)
            table.columns[2].width = Cm(3)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(5)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
            ]
            table_header = [
                "专业",
                "发版时间",
                "次数",
                "工作内容",
                "异常情况处理",
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not monthly_summary_release_data:
                    break
                data = monthly_summary_release_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 问题处理
    def slide_7(self, monthly_summary_problem_data: list):
        """问题处理"""
        # 表格行数
        row = 9
        # 表格页数初始化
        pages = 0
        if len(monthly_summary_problem_data) == 0:
            pages = 1
        elif len(monthly_summary_problem_data) % (row - 1):
            pages = len(monthly_summary_problem_data) // (row - 1) + 1
        else:
            pages = len(monthly_summary_problem_data) // (row - 1)

        for page in range(1, pages + 1):
            # 添加一张幻灯片
            self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])
            # 调整幻灯片位置
            self.insert_slide(self._current_slide_idx + 1)
            self._current_slide_idx += 1

            # 当前slide
            slide = self._prs.slides[self._current_slide_idx]  # type:pptx.slide.Slide

            # 幻灯片标题
            slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
            p = slidePlaceholder.text_frame.paragraphs[0]
            p.text = "2、运维工作进展-问题处理"
            set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

            # 创建table
            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
            x, y, cx, cy = Cm(0.5), Cm(2.5), Cm(1), Cm(0.5)
            table = slide.shapes.add_table(row, 5, x, y, cx, cy).table
            table.columns[0].width = Cm(2)
            table.columns[1].width = Cm(5)
            table.columns[2].width = Cm(5)
            table.columns[3].width = Cm(5)
            table.columns[4].width = Cm(5)
            first_row_cells = [
                table.cell(0, 0),
                table.cell(0, 1),
                table.cell(0, 2),
                table.cell(0, 3),
                table.cell(0, 4),
            ]
            table_header = [
                "专业",
                "问题详细描述",
                "处理结果",
                "原因分析",
                "后续建议",
            ]

            # 初始化所有cell的font格式
            for cell in table.iter_cells():
                tf = cell.text_frame
                tfp = cell.text_frame.paragraphs[0]
                set_text_frame_paragraph_format(tfp,
                                                font_size=11,
                                                align=PP_ALIGN.LEFT)

            for index in range(len(first_row_cells)):
                p = first_row_cells[index].text_frame.paragraphs[0]
                p.text = table_header[index]
                set_text_frame_paragraph_format(p,
                                                font_size=16,
                                                font_bold=True,
                                                align=PP_ALIGN.CENTER,
                                                font_color="FFFFFF")

            for row_idx in range(1, len(table.rows)):
                if not monthly_summary_problem_data:
                    break
                data = monthly_summary_problem_data.pop(0)
                for col_idx in range(len(table.columns)):
                    tf = table.cell(row_idx, col_idx).text_frame
                    tfp = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
                    # 自动换行
                    tf.word_wrap = True
                    tfp.text = data[col_idx]
                    set_text_frame_paragraph_format(tfp,
                                                    font_size=11,
                                                    align=PP_ALIGN.LEFT)

    # 10. 运行情况分析
    def slide_8(self, monthly_summary_analyse_data: list, months: list):
        # 解析数据
        data_length = len(monthly_summary_analyse_data)
        # cluster_names = ["FCP业务集群", "内网微服务区", "内网后台区", "外网微服务区", "外网后台区"]
        cluster_names = ["FCP业务集群"]
        cpu_allocated = {}
        memory_allocated = {}
        cluster_table_data = {}
        cluster_summary = {}
        for item in cluster_names:
            cpu_allocated[item] = []
            memory_allocated[item] = []
            cluster_table_data[item] = []

        for index in range(data_length):
            for key in monthly_summary_analyse_data[index].keys():
                cpu_allocated[key].append(monthly_summary_analyse_data[index][key]["cpu"][0])
                memory_allocated[key].append(monthly_summary_analyse_data[index][key]["memory"][0])
                if index == (data_length - 1):
                    cluster_table_data[key] = monthly_summary_analyse_data[index][key]["table_data"]
                    cluster_summary[key] = monthly_summary_analyse_data[index][key]["summary"]

        # 添加一张幻灯片
        self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

        # 当前slide
        slide = self._prs.slides[self._current_slide_idx + 5]  # type:pptx.slide.Slide

        # 幻灯片标题
        slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
        p = slidePlaceholder.text_frame.paragraphs[0]
        p.text = "四、慧企平台"
        set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

        txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2), Cm(13), Cm(1))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # # 自动换行
        # tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "运维情况分析---集群资源可分配容量"
        set_text_frame_paragraph_format(tfp, font_size=16)

        # cpu资源使用情况
        create_column_chart(slide=slide,
                            categories=cluster_names,
                            series={
                                "已分配": [monthly_summary_analyse_data[-1][key]["cpu"][0] for key in cluster_names],
                                "未分配": [1 - monthly_summary_analyse_data[-1][key]["cpu"][0] for key in
                                           cluster_names]},
                            position_and_size=[Cm(1.1), Cm(4.5), Cm(11.5), Cm(8.5)],
                            chart_title="CPU资源使用情况",
                            chart_has_legend=True,
                            column_chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                            show_data_labels=True,
                            y_axis_show_percent=True)
        # 内存资源使用情况
        create_column_chart(slide=slide,
                            categories=cluster_names,
                            series={
                                "已分配": [monthly_summary_analyse_data[-1][key]["memory"][0] for key in cluster_names],
                                "未分配": [1 - monthly_summary_analyse_data[-1][key]["memory"][0] for key in
                                           cluster_names]},
                            position_and_size=[Cm(13.6), Cm(4.5), Cm(11.5), Cm(8.5)],
                            chart_title="内存资源使用情况",
                            chart_has_legend=True,
                            column_chart_type=XL_CHART_TYPE.COLUMN_STACKED,
                            show_data_labels=True,
                            y_axis_show_percent=True)

        for cluster_name in cluster_names:
            table_data = cluster_table_data[cluster_name]

            # 定义每个集群（运行情况分析）ppt页数
            # 每页最多存60行数据
            pages = 0
            if len(table_data) % 60:
                pages = len(table_data) // 60 + 1
            else:
                pages = len(table_data) // 60

            for page in range(1, pages + 1):
                # 添加一张幻灯片
                slide = self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

                # 幻灯片标题
                slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
                p = slidePlaceholder.text_frame.paragraphs[0]
                p.text = "四、慧企平台"
                set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

                txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2.5), Cm(8), Cm(1))
                # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                tf = txbox.text_frame
                # # 自动换行
                # tf.word_wrap = True
                tfp = tf.paragraphs[0]
                tfp.text = cluster_name + "资源使用情况"
                set_text_frame_paragraph_format(tfp, font_size=16)

                # 资源使用情况折线图
                create_line_chart(slide=slide,
                                  categories=months[:data_length],
                                  series={"CPU使用率": cpu_allocated[cluster_name],
                                          "内存使用率": memory_allocated[cluster_name]},
                                  position_and_size=[Cm(0.2), Cm(4), Cm(6), Cm(7.5)],
                                  chart_title="资源使用率",
                                  chart_has_legend=True,
                                  y_axis_show_percent=True)
                if page == 1 and pages == page:
                    if len(table_data) <= 30:
                        # table_1
                        # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_summary[cluster_name]
                        set_text_frame_paragraph_format(tfp)

                    elif len(table_data) > 30:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # 右下插入文本框
                        txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                        tf = txbox.text_frame
                        # 自动换行
                        tf.word_wrap = True
                        tfp = tf.paragraphs[0]
                        tfp.text = cluster_summary[cluster_name]
                        set_text_frame_paragraph_format(tfp)
                elif page == 1:
                    # table_1
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.rows[0].heitht = Cm(0.3)
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                        # 处理集群的第一行数据，即（整个集群的表头["项目名称", "服务名称", "实例数"]）
                        for cell in table.rows[0].cells:
                            set_text_frame_paragraph_format(cell.text_frame.paragraphs[0],
                                                            font_size=7,
                                                            font_bold=True,
                                                            font_color="FFFFFF")

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="F6E7E7")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="CC0000")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])

                    # table_2
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                    x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                    table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                    table.columns[0].width = Cm(3.5)
                    table.columns[1].width = Cm(4)
                    table.columns[2].width = Cm(1)
                    # 取消表格第一行特殊格式
                    table.first_row = False
                    index = 0
                    cells = list(table.iter_cells())
                    count = len(table_data)
                    for i in range(count):
                        if index == len(cells):
                            break
                        for j in table_data.pop(0):
                            if index < len(cells):
                                cells[index].text = str(j)
                                set_cell_format(cells[index],
                                                margin_left=0,
                                                margin_right=0,
                                                margin_top=0,
                                                margin_bottom=0)
                                set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                index += 1

                    # 合并单元格
                    merge_cells_index = []
                    for index in range(len(table.rows)):
                        if table.rows[index].cells[0].text:
                            merge_cells_index.append(index)
                    # 处理首尾行是空字符的情况
                    if merge_cells_index[0] != 0:
                        merge_cells_index.insert(0, 0)
                    if merge_cells_index[-1] + 1 != len(table.rows):
                        merge_cells_index.append(len(table.rows) - 1)
                    for index in range(len(merge_cells_index)):
                        # 统一（合并单元格）颜色
                        set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                        cell_background_color="F6E7E7")
                        if table.rows[0].cells[0].text == "项目名称":
                            set_cell_format(table.rows[0].cells[0],
                                            cell_background_color="CC0000")
                        # 合并首列单元格
                        if index + 1 == len(merge_cells_index):
                            break
                        # 当最后一行的第一列为空时
                        if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                len(merge_cells_index) - 2):
                            table.rows[merge_cells_index[-2]].cells[0].merge(
                                table.rows[merge_cells_index[-1]].cells[0])
                        # 不合并自己（合并索引相邻）
                        elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                            table.rows[merge_cells_index[index]].cells[0].merge(
                                table.rows[merge_cells_index[index + 1] - 1].cells[0])
                else:
                    if page == pages:
                        if len(table_data) <= 30:
                            # table_1
                            # 根据数据长度设置表格的行列数（同时通过行高*行数，计算表格高度，并与模板对比，做表格当前页分割，或换页分割）！！！！！
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并相邻索引）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_summary[cluster_name]
                            set_text_frame_paragraph_format(tfp)
                        elif len(table_data) > 30:
                            # table_1
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)

                            # 取消表格第一行特殊格式
                            table.first_row = False

                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # table_2
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                            x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                            table = slide.shapes.add_table(len(table_data), 3, x, y, cx, cy).table
                            table.columns[0].width = Cm(3.5)
                            table.columns[1].width = Cm(4)
                            table.columns[2].width = Cm(1)
                            # 取消表格第一行特殊格式
                            table.first_row = False
                            index = 0
                            cells = list(table.iter_cells())
                            count = len(table_data)
                            for i in range(count):
                                if index == len(cells):
                                    break
                                for j in table_data.pop(0):
                                    if index < len(cells):
                                        cells[index].text = str(j)
                                        set_cell_format(cells[index],
                                                        margin_left=0,
                                                        margin_right=0,
                                                        margin_top=0,
                                                        margin_bottom=0)
                                        set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0],
                                                                        font_size=7)
                                        index += 1

                            # 合并单元格
                            merge_cells_index = []
                            for index in range(len(table.rows)):
                                if table.rows[index].cells[0].text:
                                    merge_cells_index.append(index)
                            # 处理首尾行是空字符的情况
                            if merge_cells_index[0] != 0:
                                merge_cells_index.insert(0, 0)
                            if merge_cells_index[-1] + 1 != len(table.rows):
                                merge_cells_index.append(len(table.rows) - 1)
                            for index in range(len(merge_cells_index)):
                                # 统一（合并单元格）颜色
                                set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                                cell_background_color="F6E7E7")
                                if table.rows[0].cells[0].text == "项目名称":
                                    set_cell_format(table.rows[0].cells[0],
                                                    cell_background_color="CC0000")
                                # 合并首列单元格
                                if index + 1 == len(merge_cells_index):
                                    break
                                # 当最后一行的第一列为空时
                                if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                        len(merge_cells_index) - 2):
                                    table.rows[merge_cells_index[-2]].cells[0].merge(
                                        table.rows[merge_cells_index[-1]].cells[0])
                                # 不合并自己（合并索引相邻）
                                elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                    table.rows[merge_cells_index[index]].cells[0].merge(
                                        table.rows[merge_cells_index[index + 1] - 1].cells[0])

                            # 右下插入文本框
                            txbox = slide.shapes.add_textbox(Cm(16.5), Cm(12.5), Cm(8), Cm(1))
                            # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
                            tf = txbox.text_frame
                            # 自动换行
                            tf.word_wrap = True
                            tfp = tf.paragraphs[0]
                            tfp.text = cluster_summary[cluster_name]
                            set_text_frame_paragraph_format(tfp)
                    else:
                        # table_1
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(7.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.rows[0].heitht = Cm(0.3)
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

                        # table_2
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        # x: 左边距，y: 上边距, cx: 单元格宽度, cy: 单元格高度
                        x, y, cx, cy = Cm(16.5), Cm(2.2), Cm(1), Cm(0.2)
                        table = slide.shapes.add_table(30, 3, x, y, cx, cy).table
                        table.columns[0].width = Cm(3.5)
                        table.columns[1].width = Cm(4)
                        table.columns[2].width = Cm(1)
                        # 取消表格第一行特殊格式
                        table.first_row = False
                        index = 0
                        cells = list(table.iter_cells())
                        count = len(table_data)
                        for i in range(count):
                            if index == len(cells):
                                break
                            for j in table_data.pop(0):
                                if index < len(cells):
                                    cells[index].text = str(j)
                                    set_cell_format(cells[index],
                                                    margin_left=0,
                                                    margin_right=0,
                                                    margin_top=0,
                                                    margin_bottom=0)
                                    set_text_frame_paragraph_format(cells[index].text_frame.paragraphs[0], font_size=7)
                                    index += 1

                        # 合并单元格
                        merge_cells_index = []
                        for index in range(len(table.rows)):
                            if table.rows[index].cells[0].text:
                                merge_cells_index.append(index)
                        # 处理首尾行是空字符的情况
                        if merge_cells_index[0] != 0:
                            merge_cells_index.insert(0, 0)
                        if merge_cells_index[-1] + 1 != len(table.rows):
                            merge_cells_index.append(len(table.rows) - 1)
                        for index in range(len(merge_cells_index)):
                            # 统一（合并单元格）颜色
                            set_cell_format(table.rows[merge_cells_index[index]].cells[0],
                                            cell_background_color="F6E7E7")
                            if table.rows[0].cells[0].text == "项目名称":
                                set_cell_format(table.rows[0].cells[0],
                                                cell_background_color="CC0000")
                            # 合并首列单元格
                            if index + 1 == len(merge_cells_index):
                                break
                            # 当最后一行的第一列为空时
                            if not table.rows[merge_cells_index[-1]].cells[0].text and index == (
                                    len(merge_cells_index) - 2):
                                table.rows[merge_cells_index[-2]].cells[0].merge(
                                    table.rows[merge_cells_index[-1]].cells[0])
                            # 不合并自己（合并索引相邻）
                            elif merge_cells_index[index] + 1 != merge_cells_index[index + 1]:
                                table.rows[merge_cells_index[index]].cells[0].merge(
                                    table.rows[merge_cells_index[index + 1] - 1].cells[0])

        # 添加一张幻灯片
        slide = self._prs.slides.add_slide(self._prs.slide_masters[0].slide_layouts[1])

        # 幻灯片标题
        slidePlaceholder = slide.shapes[0]  # type: pptx.shapes.placeholder.SlidePlaceholder
        p = slidePlaceholder.text_frame.paragraphs[0]
        p.text = "四、慧企平台"
        set_text_frame_paragraph_format(p, font_bold=True, font_size=24)

        txbox = slide.shapes.add_textbox(Cm(0.1), Cm(2), Cm(13), Cm(1))
        # 文本框应该有一个默认段落，直接获取，额外添加段落需要使用text_frame.add_paragraph()，清除段落用text_frame.clear()
        tf = txbox.text_frame
        # # 自动换行
        # tf.word_wrap = True
        tfp = tf.paragraphs[0]
        tfp.text = "工作中的问题或困难"
        set_text_frame_paragraph_format(tfp, font_size=16)
